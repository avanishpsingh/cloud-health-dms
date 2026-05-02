"""Generate a polished Phase 2 PowerPoint deck.

Reads `report/Phase2_PPT_Outline.md` and produces a 16:9 deck with:

* Custom title slide (navy banner, project + team).
* Section divider slides for "Live Demo", "Reflection", "Q&A".
* Content slides with a left-aligned accent bar, capped title, bulleted
  list with a teal glyph, code blocks in a fixed-width box, and pipe
  tables rendered as real PowerPoint tables.
* Footer on every slide with project name + slide number.
* Speaker notes for every slide (italic outline lines).

Run:
    .\\venv\\Scripts\\python.exe scripts\\generate_pptx.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parent.parent
SRC = REPO / "report" / "Phase2_PPT_Outline.md"
OUT = REPO / "report" / "CC_Project_PPT_Group23.pptx"

# ---------------------------------------------------------------------------
# Palette — consistent BITS-blue + accent
# ---------------------------------------------------------------------------
NAVY     = RGBColor(0x0B, 0x2E, 0x6F)
BLUE     = RGBColor(0x1E, 0x55, 0xB3)
TEAL     = RGBColor(0x10, 0x9B, 0xC4)
INK      = RGBColor(0x1A, 0x1A, 0x1A)
MUTED    = RGBColor(0x55, 0x5F, 0x6D)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFB)
ROW_ALT  = RGBColor(0xEE, 0xF2, 0xF9)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# Geometry (16:9 — 13.333 x 7.5 inches)
SLIDE_W   = Inches(13.333)
SLIDE_H   = Inches(7.5)
MARGIN    = Inches(0.55)
TITLE_TOP = Inches(0.40)
TITLE_H   = Inches(0.85)
BODY_TOP  = Inches(1.55)
BODY_H    = Inches(5.35)


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------
def parse_slides(text: str) -> list[dict]:
    blocks = re.split(r"\n---+\s*\n", text)
    slides: list[dict] = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        m = re.search(r"^##\s+(.*?)\s*$", block, re.MULTILINE)
        if not m:
            continue
        full_title = m.group(1).strip()
        title = re.sub(r"^Slide\s+\d+\s*[—–-]\s*", "", full_title).strip()
        body = block[m.end():].strip()
        slides.append({"title": title, "body": body})
    return slides


def strip_inline_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text


def parse_table(lines: list[str]) -> list[list[str]] | None:
    rows = []
    for line in lines:
        if "|" not in line:
            return None
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        rows.append(cells)
    if len(rows) >= 2 and all(re.fullmatch(r"[\-:\s]+", c) for c in rows[1]):
        rows.pop(1)
    return rows


def split_blocks(body: str):
    lines = body.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        s = line.strip()

        if s.startswith("*") and s.endswith("*") and not s.startswith("**"):
            yield ("note", s.strip("*").strip())
            i += 1
            continue

        if s.startswith("```"):
            j = i + 1
            buf = []
            while j < len(lines) and not lines[j].lstrip().startswith("```"):
                buf.append(lines[j])
                j += 1
            yield ("code", "\n".join(buf))
            i = j + 1
            continue

        if s.startswith(("- ", "* ")):
            buf = []
            while i < len(lines) and lines[i].strip().startswith(("- ", "* ")):
                buf.append(lines[i].strip()[2:])
                i += 1
            yield ("bullets", buf)
            continue

        if s.startswith("|"):
            buf = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                buf.append(lines[i])
                i += 1
            data = parse_table(buf)
            if data:
                yield ("table", data)
            else:
                yield ("para", " ".join(b.strip() for b in buf))
            continue

        if not s:
            i += 1
            continue

        buf = [s]
        i += 1
        while i < len(lines):
            nxt = lines[i].strip()
            if not nxt:
                break
            if nxt.startswith(("- ", "* ", "|", "```", "*")):
                break
            buf.append(nxt)
            i += 1
        yield ("para", " ".join(buf))


# ---------------------------------------------------------------------------
# Drawing primitives
# ---------------------------------------------------------------------------
def add_filled_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_footer(slide, idx: int, total: int):
    add_text(slide, MARGIN, Inches(7.05), Inches(8), Inches(0.3),
             "Cloud-Native Healthcare DMS  ·  BITS Pilani  ·  Group 23",
             size=9, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title_block(slide, title: str):
    add_filled_rect(slide, MARGIN, TITLE_TOP, Inches(0.18), TITLE_H, TEAL)
    tb = slide.shapes.add_textbox(
        Inches(0.95), TITLE_TOP, SLIDE_W - Inches(1.5), TITLE_H
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = strip_inline_md(title)
    r.font.name = "Calibri"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY
    add_filled_rect(slide, Inches(0.95), Inches(1.30), Inches(1.2),
                    Emu(38100), TEAL)


# ---------------------------------------------------------------------------
# Title / divider slides
# ---------------------------------------------------------------------------
def render_title_slide(prs, _title: str, subtitle_lines: list[str]):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(4.4), NAVY)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(4.35),
                                    SLIDE_W, Inches(0.12))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()

    add_text(slide, MARGIN, Inches(1.15), SLIDE_W - 2 * MARGIN, Inches(1.2),
             "Cloud-Native Healthcare DMS", size=44, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, MARGIN, Inches(2.15), SLIDE_W - 2 * MARGIN, Inches(0.7),
             "Phase 2 — Implementation & Demonstration", size=24,
             color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, MARGIN, Inches(2.95), SLIDE_W - 2 * MARGIN, Inches(0.6),
             "BITS Pilani  ·  Cloud Computing (CSIZG527)  ·  Group 23",
             size=16, color=WHITE, align=PP_ALIGN.LEFT)

    y = Inches(4.85)
    for line in subtitle_lines:
        add_text(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.45),
                 strip_inline_md(line), size=18, color=INK)
        y += Inches(0.45)


def render_divider_slide(prs, title: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_filled_rect(slide, Inches(0), Inches(3.5), SLIDE_W, Emu(38100), TEAL)
    add_text(slide, MARGIN, Inches(2.7), SLIDE_W - 2 * MARGIN, Inches(0.6),
             "Section", size=18, color=TEAL, align=PP_ALIGN.LEFT, bold=True)
    add_text(slide, MARGIN, Inches(3.7), SLIDE_W - 2 * MARGIN, Inches(1.5),
             strip_inline_md(title), size=44, bold=True, color=WHITE)
    return slide


# ---------------------------------------------------------------------------
# Content rendering
# ---------------------------------------------------------------------------
def render_content_slide(prs, title: str, body: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, title)

    body_x = MARGIN
    body_y = BODY_TOP
    body_w = SLIDE_W - 2 * MARGIN
    cursor_y = body_y

    notes_buf: list[str] = []

    for kind, content in split_blocks(body):
        if cursor_y > body_y + BODY_H - Inches(0.4):
            # Out of room — drop remaining blocks into speaker notes
            if kind == "para":
                notes_buf.append(strip_inline_md(content))
            elif kind == "bullets":
                notes_buf.extend(strip_inline_md(b) for b in content)
            elif kind == "code":
                notes_buf.append(content)
            elif kind == "note":
                notes_buf.append(content)
            continue

        if kind == "note":
            notes_buf.append(content)
            continue

        if kind == "para":
            txt = strip_inline_md(content)
            tb = slide.shapes.add_textbox(body_x, cursor_y, body_w, Inches(0.6))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = txt
            r.font.name = "Calibri"; r.font.size = Pt(18); r.font.color.rgb = INK
            cursor_y += Inches(0.55)

        elif kind == "bullets":
            n = len(content)
            tb = slide.shapes.add_textbox(body_x, cursor_y, body_w,
                                          Inches(min(0.5 * n + 0.2, 5.0)))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_top = Emu(0)
            for i, item in enumerate(content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(4)
                r1 = p.add_run(); r1.text = "▸  "
                r1.font.name = "Calibri"; r1.font.size = Pt(18)
                r1.font.color.rgb = TEAL; r1.font.bold = True
                r2 = p.add_run(); r2.text = strip_inline_md(item)
                r2.font.name = "Calibri"; r2.font.size = Pt(17)
                r2.font.color.rgb = INK
            cursor_y += Inches(0.42 * n + 0.25)

        elif kind == "code":
            lines = content.splitlines() or [""]
            box_h = Inches(min(0.30 * len(lines) + 0.30, 4.5))
            add_filled_rect(slide, body_x, cursor_y, body_w, box_h, LIGHT_BG)
            tb = slide.shapes.add_textbox(
                body_x + Inches(0.15), cursor_y + Inches(0.10),
                body_w - Inches(0.3), box_h - Inches(0.20)
            )
            tf = tb.text_frame; tf.word_wrap = False
            tf.margin_left = Emu(0); tf.margin_top = Emu(0)
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run(); r.text = line
                r.font.name = "Consolas"; r.font.size = Pt(11)
                r.font.color.rgb = MUTED
            cursor_y += box_h + Inches(0.15)

        elif kind == "table":
            data = content
            rows = len(data); cols = len(data[0])
            row_h = Inches(0.40)
            tbl_h = row_h * rows
            shape = slide.shapes.add_table(rows, cols, body_x, cursor_y,
                                           body_w, tbl_h)
            tbl = shape.table
            if cols >= 2:
                tbl.columns[0].width = int(body_w * 0.30)
                rest = body_w - tbl.columns[0].width
                for c in range(1, cols):
                    tbl.columns[c].width = int(rest / (cols - 1))
            for r_idx, row in enumerate(data):
                for c_idx, val in enumerate(row):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.margin_left = Inches(0.08)
                    cell.margin_right = Inches(0.08)
                    cell.margin_top = Inches(0.04)
                    cell.margin_bottom = Inches(0.04)
                    cell.fill.solid()
                    if r_idx == 0:
                        cell.fill.fore_color.rgb = NAVY
                    else:
                        cell.fill.fore_color.rgb = ROW_ALT if r_idx % 2 == 0 else WHITE
                    tf = cell.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = strip_inline_md(val)
                    run.font.name = "Calibri"
                    run.font.size = Pt(12)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = INK
            cursor_y += tbl_h + Inches(0.20)

    if notes_buf:
        notes = slide.notes_slide.notes_text_frame
        notes.text = "\n".join(notes_buf)

    return slide


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------
DIVIDER_KEYWORDS = ("Live Demo", "Q&A")


def main() -> None:
    text = SRC.read_text(encoding="utf-8")
    slides_md = parse_slides(text)
    if not slides_md:
        raise SystemExit("No slides parsed from outline")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title = slides_md[0]
    subtitle = [
        line.strip() for line in title["body"].splitlines()
        if line.strip() and not line.strip().startswith("*")
    ]
    render_title_slide(prs, title["title"], subtitle)

    rest = slides_md[1:]
    total = len(rest) + 1
    for idx, s in enumerate(rest, start=2):
        if any(kw in s["title"] for kw in DIVIDER_KEYWORDS):
            render_divider_slide(prs, s["title"])
        else:
            render_content_slide(prs, s["title"], s["body"])
        add_footer(prs.slides[-1], idx, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"✓ Wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
